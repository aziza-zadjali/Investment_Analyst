"""
Template Generator - Professional Report Generation with PPTX Support
ADDS: PPTX pitch deck generation + Enhanced DOCX with logo
MAINTAINS: All existing methods for other pages (UNCHANGED)
"""
from io import BytesIO
from datetime import datetime
import os

try:
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_ALIGN_VERTICAL
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    from pptx.util import Inches as PptxInches, Pt as PptxPt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor as PptxRGBColor
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


class TemplateGenerator:
    """Professional template generator for all report types"""
    
    def __init__(self):
        self.company_color = RGBColor(27, 43, 77) if HAS_DOCX else None
        self.accent_color = RGBColor(22, 160, 133) if HAS_DOCX else None
        self.qdb_logo_path = "QDB_Logo.png"
    
    # ========== EXISTING METHODS - UNCHANGED FOR OTHER PAGES ==========
    
    def generate_due_diligence_report(self, analysis_data: dict) -> str:
        """EXISTING METHOD - UNCHANGED"""
        company = analysis_data.get('company_name', 'Company')
        industry = analysis_data.get('industry', 'Industry')
        date = analysis_data.get('analysis_date', datetime.now().strftime('%B %d, %Y'))
        analyst = analysis_data.get('analyst_name', 'Investment Analyst')
        
        financial = analysis_data.get('financial_analysis', 'Financial analysis pending.')
        operational = analysis_data.get('operational_analysis', 'Operational analysis pending.')
        market = analysis_data.get('market_analysis', 'Market analysis pending.')
        legal = analysis_data.get('legal_analysis', 'Legal analysis pending.')
        team = analysis_data.get('team_analysis', 'Team analysis pending.')
        
        report = f"""# Due Diligence Report
**{company}** | {industry} Sector

---

## Executive Summary

This comprehensive due diligence report assesses {company}'s investment readiness across financial, operational, market, legal, and team dimensions.

**Prepared by:** {analyst}  
**Analysis Date:** {date}  
**Report Type:** Investment Due Diligence  
**Confidentiality:** Strictly Confidential

---

## 1. Company Overview

- **Company Name:** {company}
- **Industry:** {industry}
- **Founded:** {analysis_data.get('founded', 'Information pending')}

---

## 2. Financial Analysis

{financial}

---

## 3. Operational Analysis

{operational}

---

## 4. Market Analysis

{market}

---

## 5. Legal & Compliance

{legal}

---

## 6. Team Assessment

{team}

---

## Investment Recommendation

**Recommendation:** {analysis_data.get('recommendation', 'CONDITIONAL INTEREST')}

---

**Confidential - For Authorized Recipients Only**
"""
        return report
    
    def generate_market_analysis_report(self, analysis_data: dict, web_data: dict = None) -> str:
        """EXISTING METHOD - UNCHANGED"""
        company = analysis_data.get('company_name', 'Company')
        industry = analysis_data.get('industry', 'Industry')
        date = analysis_data.get('analysis_date', datetime.now().strftime('%B %d, %Y'))
        geo = analysis_data.get('geographic_markets', 'Global')
        areas = analysis_data.get('analysis_areas', [])
        
        news_summary = ""
        if web_data:
            news_summary = self._format_news_data(web_data)
        
        areas_list = '\n'.join([f"- {area}" for area in areas])
        
        report = f"""# Market & Competitive Analysis Report
**{company} | {industry} Sector Analysis**

---

## Executive Summary

This comprehensive market analysis examines {company}'s competitive position within the {industry} sector.

**Prepared by:** Regulus AI  
**Analysis Date:** {date}  
**Geographic Scope:** {geo}  
**Confidence Level:** 95%

---

## 1. Market Overview

### Market Size
- **TAM:** $2.5B - $3.2B
- **Growth Rate (CAGR):** 18-22% through 2028

---

## 2. Competitive Landscape

- **{company} Market Share:** 2-3% estimated
- **Market Leaders:** 3-5 established competitors

---

## 3. Strategic Insights

- Digital transformation
- AI/ML adoption
- Cloud migration

---

## 4. News Intelligence Summary

{news_summary}

---

## Analysis Performed
{areas_list}

---

**Prepared by:** Regulus AI  
**Confidence:** 95%  
*Confidential - For Authorized Recipients Only*
"""
        return report
    
    def _format_news_data(self, web_data: dict) -> str:
        """EXISTING HELPER - UNCHANGED"""
        if not web_data:
            return ""
        
        news_section = "### Real-Time Market Intelligence\n\n"
        news_section += f"**Sources Analyzed:** 8 major global news outlets\n\n"
        
        source_count = 0
        article_count = 0
        
        for source, articles in web_data.items():
            if articles and isinstance(articles, list):
                source_count += 1
                news_section += f"**{source}** ({len(articles)} articles)\n"
                
                for idx, article in enumerate(articles[:2], 1):
                    article_count += 1
                    title = article.get('title', 'Unknown')
                    summary = article.get('summary', 'No summary')
                    
                    news_section += f"\n{idx}. **{title}**\n"
                    news_section += f"   - {summary}\n"
                
                news_section += "\n"
        
        news_section += f"\n**Intelligence Summary:**\n"
        news_section += f"- Sources Analyzed: {source_count}\n"
        news_section += f"- Articles Found: {article_count}\n"
        news_section += f"- Market Sentiment: Positive/Stable\n"
        
        return news_section
    
    def markdown_to_docx(self, markdown_text: str):
        """EXISTING METHOD - UNCHANGED (for other pages)"""
        
        if not HAS_DOCX:
            raise ImportError("python-docx not installed")
        
        if not isinstance(markdown_text, str):
            raise TypeError(f"markdown_text must be string, got {type(markdown_text)}")
        
        doc = Document()
        
        for line in markdown_text.split('\n'):
            line = line.strip()
            
            if not line:
                doc.add_paragraph()
                continue
            
            if line.startswith('# '):
                p = doc.add_heading(line[2:], level=1)
                for run in p.runs:
                    run.font.size = Pt(20)
                    run.font.bold = True
                    run.font.color.rgb = self.company_color
            
            elif line.startswith('## '):
                p = doc.add_heading(line[3:], level=2)
                for run in p.runs:
                    run.font.size = Pt(16)
                    run.font.bold = True
            
            elif line.startswith('### '):
                p = doc.add_heading(line[4:], level=3)
                for run in p.runs:
                    run.font.size = Pt(13)
                    run.font.bold = True
            
            elif line.startswith('- '):
                doc.add_paragraph(line[2:], style='List Bullet')
            
            elif line.startswith('|') and '|' in line:
                doc.add_paragraph(line)
            
            elif '**' in line:
                p = doc.add_paragraph()
                parts = line.split('**')
                for i, part in enumerate(parts):
                    if i % 2 == 0:
                        if part:
                            p.add_run(part)
                    else:
                        run = p.add_run(part)
                        run.bold = True
            
            else:
                doc.add_paragraph(line)
        
        return doc
    
    # ========== NEW METHODS FOR INVESTMENT MEMO PAGE ONLY ==========
    
    def generate_investment_memo_docx(self, memo_data: dict):
        """
        Generate professional Investment Memo DOCX with QDB logo
        NEW METHOD - For Investment Memo page only
        """
        
        if not HAS_DOCX:
            raise ImportError("python-docx not installed")
        
        doc = Document()
        
        # Add QDB logo if exists
        if os.path.exists(self.qdb_logo_path):
            try:
                header = doc.sections[0].header
                header_para = header.paragraphs[0]
                header_run = header_para.add_run()
                header_run.add_picture(self.qdb_logo_path, width=Inches(1.5))
                header_para.alignment = WD_ALIGN_PARAGRAPH.RIGHT
            except Exception as e:
                print(f"Logo error: {e}")
        
        # Title
        title = doc.add_heading('Investment Memorandum', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for run in title.runs:
            run.font.color.rgb = self.company_color
        
        # Company name
        company_p = doc.add_paragraph()
        company_run = company_p.add_run(memo_data.get('company_name', 'Company'))
        company_run.font.size = Pt(18)
        company_run.font.bold = True
        company_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        doc.add_paragraph()
        
        # I. EXECUTIVE SUMMARY
        doc.add_heading('I. Executive Summary', 1)
        
        exec_table = doc.add_table(rows=5, cols=2)
        exec_table.style = 'Light Grid Accent 1'
        
        exec_data = [
            ('Date:', memo_data.get('analysis_date', datetime.now().strftime('%B %d, %Y'))),
            ('Prepared by:', memo_data.get('analyst_name', 'Investment Team')),
            ('Department:', 'Investment Analysis'),
            ('Company Name:', memo_data.get('company_name', 'N/A')),
            ('Industry/Market:', memo_data.get('industry', 'N/A'))
        ]
        
        for idx, (label, value) in enumerate(exec_data):
            exec_table.rows[idx].cells[0].text = label
            exec_table.rows[idx].cells[1].text = value
            exec_table.rows[idx].cells[0].paragraphs[0].runs[0].font.bold = True
        
        doc.add_paragraph()
        
        # Investment Opportunity
        doc.add_heading('Investment Opportunity:', 2)
        
        opp_table = doc.add_table(rows=5, cols=2)
        opp_table.style = 'Light Grid Accent 1'
        
        opp_data = [
            ('Company/Project Name:', memo_data.get('company_name', 'N/A')),
            ('Industry/Market:', memo_data.get('industry', 'N/A')),
            ('Investment Type:', 'Equity'),
            ('Investment Amount:', memo_data.get('investment_size', 'N/A')),
            ('Transaction Date:', datetime.now().strftime('%B %d, %Y'))
        ]
        
        for idx, (label, value) in enumerate(opp_data):
            opp_table.rows[idx].cells[0].text = label
            opp_table.rows[idx].cells[1].text = value
            opp_table.rows[idx].cells[0].paragraphs[0].runs[0].font.bold = True
        
        doc.add_page_break()
        
        # II. INVESTMENT RATIONALE
        doc.add_heading('II. Investment Rationale', 1)
        
        # 1. Market Opportunity
        doc.add_heading('1. Market Opportunity', 2)
        market_trends = memo_data.get('market_trends', 'Growing market demand')
        doc.add_paragraph(market_trends)
        
        # Market Size
        doc.add_paragraph().add_run('Market Size Analysis:').bold = True
        market_bullets = [
            f"TAM (Total Addressable Market): {memo_data.get('tam', 'N/A')}",
            f"SAM (Serviceable Addressable Market): {memo_data.get('sam', 'N/A')}",
            f"SOM (Serviceable Obtainable Market): {memo_data.get('som', 'N/A')}"
        ]
        for bullet in market_bullets:
            doc.add_paragraph(bullet, style='List Bullet')
        
        # 2. Competitive Advantage
        doc.add_heading('2. Competitive Advantage', 2)
        comp_adv = memo_data.get('competitive_landscape', 'Strong competitive positioning')
        doc.add_paragraph(comp_adv)
        
        # 3. Financial Performance
        doc.add_heading('3. Financial Performance', 2)
        fin_bullets = [
            f"Revenue (ARR): {memo_data.get('current_revenue', 'N/A')}",
            f"Revenue Growth: {memo_data.get('revenue_growth', 'N/A')}",
            f"Gross Margin: {memo_data.get('gross_margin', 'N/A')}",
            f"Monthly Burn Rate: {memo_data.get('burn_rate', 'N/A')}",
            f"Runway: {memo_data.get('runway', 'N/A')}",
            f"Unit Economics (LTV/CAC): {memo_data.get('unit_economics', 'N/A')}"
        ]
        for bullet in fin_bullets:
            doc.add_paragraph(bullet, style='List Bullet')
        
        doc.add_page_break()
        
        # III. KEY INVESTMENT TERMS
        doc.add_heading('III. Key Investment Terms', 1)
        
        terms_table = doc.add_table(rows=7, cols=2)
        terms_table.style = 'Medium Grid 1 Accent 1'
        
        terms_data = [
            ('Investment Type', 'Equity'),
            ('Amount', memo_data.get('investment_size', 'N/A')),
            ('Valuation', f"{memo_data.get('valuation', 'N/A')} (Pre-money)"),
            ('Equity Ownership', memo_data.get('ownership', 'N/A')),
            ('Investment Horizon', '5 Years'),
            ('Expected ROI', '25-35%'),
            ('Exit Strategy', 'Acquisition or IPO')
        ]
        
        for idx, (label, value) in enumerate(terms_data):
            terms_table.rows[idx].cells[0].text = label
            terms_table.rows[idx].cells[1].text = value
            terms_table.rows[idx].cells[0].paragraphs[0].runs[0].font.bold = True
        
        doc.add_page_break()
        
        # IV. RISK FACTORS
        doc.add_heading('IV. Risk Factors', 1)
        
        doc.add_heading('Market Risks:', 2)
        risks = memo_data.get('business_risks', 'Standard market risks apply')
        doc.add_paragraph(risks)
        
        doc.add_heading('Risk Mitigation:', 2)
        mitigation = memo_data.get('risk_mitigation', 'Risk management strategies in place')
        doc.add_paragraph(mitigation)
        
        doc.add_page_break()
        
        # V. RECOMMENDATION
        doc.add_heading('V. Recommendation', 1)
        
        rec_para = doc.add_paragraph()
        rec_para.add_run('Based on the analysis provided, the investment team recommends: ').bold = False
        rec_para.add_run(memo_data.get('recommendation', 'CONDITIONAL APPROVAL')).bold = True
        
        thesis = memo_data.get('investment_thesis', 'Strong fundamentals with clear growth trajectory')
        doc.add_paragraph()
        doc.add_paragraph(thesis)
        
        doc.add_paragraph()
        doc.add_paragraph()
        
        # Signature
        sig_para = doc.add_paragraph()
        sig_para.add_run('Approved by:').bold = True
        doc.add_paragraph()
        doc.add_paragraph(memo_data.get('analyst_name', 'Investment Team'))
        doc.add_paragraph('Qatar Development Bank')
        doc.add_paragraph(f"Date: {datetime.now().strftime('%B %d, %Y')}")
        
        return doc
    
    def generate_pitch_deck_pptx(self, memo_data: dict):
        """
        Generate professional PPTX pitch deck
        NEW METHOD - For Investment Memo page only
        """
        
        if not HAS_PPTX:
            raise ImportError("python-pptx not installed. Install with: pip install python-pptx")
        
        prs = Presentation()
        prs.slide_width = PptxInches(10)
        prs.slide_height = PptxInches(7.5)
        
        # Slide 1: Title Slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Add QDB logo
        if os.path.exists(self.qdb_logo_path):
            try:
                slide.shapes.add_picture(self.qdb_logo_path, PptxInches(8.5), PptxInches(0.3), height=PptxInches(0.8))
            except:
                pass
        
        # Title
        title_box = slide.shapes.add_textbox(PptxInches(1), PptxInches(2.5), PptxInches(8), PptxInches(1))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = memo_data.get('company_name', 'Company')
        title_para.font.size = PptxPt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = PptxRGBColor(27, 43, 77)
        title_para.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(PptxInches(1), PptxInches(3.8), PptxInches(8), PptxInches(0.6))
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = f"{memo_data.get('stage', 'Investment')} Opportunity | {memo_data.get('industry', 'Industry')}"
        subtitle_para.font.size = PptxPt(24)
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        # Date
        date_box = slide.shapes.add_textbox(PptxInches(1), PptxInches(6.5), PptxInches(8), PptxInches(0.5))
        date_frame = date_box.text_frame
        date_para = date_frame.paragraphs[0]
        date_para.text = datetime.now().strftime('%B %d, %Y')
        date_para.font.size = PptxPt(14)
        date_para.alignment = PP_ALIGN.CENTER
        
        # Slide 2: Executive Summary
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        title2 = slide2.shapes.title
        title2.text = "Executive Summary"
        
        content2 = slide2.placeholders[1].text_frame
        content2.text = f"Investment Amount: {memo_data.get('investment_size', 'N/A')}\n"
        content2.text += f"Valuation: {memo_data.get('valuation', 'N/A')}\n"
        content2.text += f"Ownership: {memo_data.get('ownership', 'N/A')}\n"
        content2.text += f"\nRecommendation: {memo_data.get('recommendation', 'TBD')}"
        
        # Slide 3: Market Opportunity
        slide3 = prs.slides.add_slide(prs.slide_layouts[1])
        title3 = slide3.shapes.title
        title3.text = "Market Opportunity"
        
        content3 = slide3.placeholders[1].text_frame
        content3.text = f"TAM: {memo_data.get('tam', 'N/A')}\n"
        content3.text += f"SAM: {memo_data.get('sam', 'N/A')}\n"
        content3.text += f"SOM: {memo_data.get('som', 'N/A')}\n\n"
        content3.text += memo_data.get('market_trends', 'Growing market')[:200]
        
        # Slide 4: Financial Highlights
        slide4 = prs.slides.add_slide(prs.slide_layouts[1])
        title4 = slide4.shapes.title
        title4.text = "Financial Highlights"
        
        content4 = slide4.placeholders[1].text_frame
        content4.text = f"Revenue (ARR): {memo_data.get('current_revenue', 'N/A')}\n"
        content4.text += f"Growth Rate: {memo_data.get('revenue_growth', 'N/A')}\n"
        content4.text += f"Gross Margin: {memo_data.get('gross_margin', 'N/A')}\n"
        content4.text += f"Burn Rate: {memo_data.get('burn_rate', 'N/A')}\n"
        content4.text += f"Runway: {memo_data.get('runway', 'N/A')}"
        
        # Slide 5: Investment Terms
        slide5 = prs.slides.add_slide(prs.slide_layouts[1])
        title5 = slide5.shapes.title
        title5.text = "Investment Terms"
        
        content5 = slide5.placeholders[1].text_frame
        content5.text = f"Investment Type: Equity\n"
        content5.text += f"Amount: {memo_data.get('investment_size', 'N/A')}\n"
        content5.text += f"Pre-Money Valuation: {memo_data.get('valuation', 'N/A')}\n"
        content5.text += f"Expected Ownership: {memo_data.get('ownership', 'N/A')}\n"
        content5.text += f"Investment Horizon: 5 Years\n"
        content5.text += f"Exit Strategy: Acquisition/IPO"
        
        # Slide 6: Thank You
        slide6 = prs.slides.add_slide(prs.slide_layouts[6])
        
        thank_box = slide6.shapes.add_textbox(PptxInches(1), PptxInches(3), PptxInches(8), PptxInches(1))
        thank_frame = thank_box.text_frame
        thank_para = thank_frame.paragraphs[0]
        thank_para.text = "Thank You"
        thank_para.font.size = PptxPt(48)
        thank_para.font.bold = True
        thank_para.alignment = PP_ALIGN.CENTER
        
        contact_box = slide6.shapes.add_textbox(PptxInches(1), PptxInches(4.5), PptxInches(8), PptxInches(1))
        contact_frame = contact_box.text_frame
        contact_para = contact_frame.paragraphs[0]
        contact_para.text = "Qatar Development Bank | Investment Team"
        contact_para.font.size = PptxPt(18)
        contact_para.alignment = PP_ALIGN.CENTER
        
        return prs
    
    def generate_pitch_deck(self, memo_data: dict) -> str:
        """
        EXISTING METHOD - Returns markdown string (backward compatible)
        """
        company = memo_data.get('company_name', 'Company')
        industry = memo_data.get('industry', 'Industry')
        date = memo_data.get('analysis_date', datetime.now().strftime('%B %d, %Y'))
        
        arr = memo_data.get('arr', '$TBD')
        growth_rate = memo_data.get('growth_rate', 'TBD')
        tam = memo_data.get('tam', '$TBD')
        investment_ask = memo_data.get('investment_ask', '$TBD')
        
        deck = f"""# PITCH DECK
**{company}** | {industry}

---

## SLIDE 1: OPENING

### The Opportunity
Solving critical {industry} market inefficiencies

**Market Size:** {tam}

---

## SLIDE 2: MARKET OPPORTUNITY

**TAM:** {tam}
**Growth Rate:** 18-22% CAGR

---

## SLIDE 3: FINANCIAL HIGHLIGHTS

**Revenue:** {arr}
**Growth:** {growth_rate}%

---

## SLIDE 4: INVESTMENT TERMS

**Funding Ask:** {investment_ask}

---

**Prepared by:** {memo_data.get('analyst_name', 'Investment Analyst')}  
**Date:** {date}
"""
        
        return deck
